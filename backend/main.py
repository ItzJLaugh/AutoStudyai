"""
AutoStudyAI Backend API
FastAPI server for processing educational content and generating study materials.
"""

import os
import re
import logging
import traceback
from uuid import uuid4

from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException, Request, Header, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from slowapi import Limiter, _rate_limit_exceeded_handler
from slowapi.util import get_remote_address
from slowapi.errors import RateLimitExceeded

from storage import InMemoryStorage
from schemas import (
    IngestRequest, IngestResponse,
    GenerateRequest, GenerateResponse,
    FlashcardRequest, FlashcardResponse,
    ChatRequest, ChatResponse
)
from services.text_processing import (
    clean_text, chunk_text,
    is_slideshow_content, extract_slideshow_content,
    format_slideshow_text, inject_image_descriptions,
    inject_page_image_descriptions
)
from services.llm import (
    generate_notes_ai, generate_study_guide,
    generate_flashcards, answer_question,
    analyze_images_for_slides
)
from routers import auth, folders, guides, stats, search, quiz, billing, nclex, exam, feedback, smart_notes
from auth_utils import get_user_id
from routers.billing import check_and_increment_usage

# Load environment variables
load_dotenv(dotenv_path=os.path.join(os.path.dirname(__file__), '.env'))

# Configure logging — never log tokens or secrets
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Rate limiter — keyed by IP address
limiter = Limiter(key_func=get_remote_address)

# Initialize FastAPI app
app = FastAPI(
    title="AutoStudyAI API",
    description="API for generating study materials from educational content",
    version="1.0.0",
    docs_url="/docs" if os.getenv("ENVIRONMENT") != "production" else None,
    redoc_url=None,
)

# Attach rate limiter
app.state.limiter = limiter
app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)

# CORS — restrict to known origins
ALLOWED_ORIGINS = [
    "http://localhost:3000",
    "http://127.0.0.1:3000",
    "https://autostudyai.online",
    "https://www.autostudyai.online",
]
if os.getenv("FRONTEND_URL"):
    ALLOWED_ORIGINS.append(os.getenv("FRONTEND_URL"))

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_origin_regex=r"chrome-extension://.*|https://.*\.vercel\.app",
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "PATCH", "DELETE", "OPTIONS"],
    allow_headers=["Authorization", "Content-Type"],
    max_age=600,
)


# Security headers middleware
@app.middleware("http")
async def add_security_headers(request: Request, call_next):
    response = await call_next(request)
    response.headers["X-Content-Type-Options"] = "nosniff"
    response.headers["X-Frame-Options"] = "DENY"
    response.headers["X-XSS-Protection"] = "1; mode=block"
    response.headers["Referrer-Policy"] = "strict-origin-when-cross-origin"
    response.headers["Cache-Control"] = "no-store"
    response.headers["Permissions-Policy"] = "camera=(), microphone=(), geolocation=()"
    # HSTS — enforce HTTPS in production
    if os.getenv("ENVIRONMENT") == "production":
        response.headers["Strict-Transport-Security"] = "max-age=31536000; includeSubDomains"
    return response


# Include Supabase-backed routers
app.include_router(auth.router)
app.include_router(folders.router)
app.include_router(guides.router)
app.include_router(stats.router)
app.include_router(search.router)
app.include_router(quiz.router)
app.include_router(nclex.router)
app.include_router(exam.router)
app.include_router(billing.router)
app.include_router(feedback.router)
app.include_router(smart_notes.router)

# Initialize storage with limits
storage = InMemoryStorage()

# === Input sanitization helpers ===
MAX_CONTENT_LENGTH = 500_000  # 500KB max content
MAX_QUESTION_LENGTH = 2_000
MAX_URL_LENGTH = 2_048


def _sanitize_text(text: str, max_length: int = MAX_CONTENT_LENGTH) -> str:
    """Enforce length limits on text input."""
    if not text:
        return ""
    return text[:max_length]


def _validate_url(url: str) -> str:
    """Validate and sanitize URL."""
    if not url:
        return ""
    url = url[:MAX_URL_LENGTH]
    if not re.match(r'^https?://', url):
        raise HTTPException(status_code=400, detail="Invalid URL")
    return url


@app.post("/extract-file-text")
@limiter.limit("20/minute")
async def extract_file_text(request: Request, file: UploadFile = None, authorization: str = Header(default="")):
    """Extract plain text from uploaded PDF, DOCX, PPTX, or TXT file."""
    import io
    try:
        get_user_id(authorization)
        if file is None:
            raise HTTPException(status_code=400, detail="No file uploaded")

        filename = (file.filename or "").lower()
        content_bytes = await file.read()

        if len(content_bytes) > 20 * 1024 * 1024:
            raise HTTPException(status_code=400, detail="File too large (max 20MB)")

        text = ""

        if filename.endswith(".pdf"):
            try:
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(content_bytes))
                parts = []
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        parts.append(extracted)
                text = "\n\n".join(parts)
            except HTTPException:
                raise
            except Exception:
                raise HTTPException(status_code=422, detail="Could not extract text from PDF. Use a text-based PDF.")

        elif filename.endswith(".docx"):
            try:
                from docx import Document
                doc = Document(io.BytesIO(content_bytes))
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            except Exception:
                raise HTTPException(status_code=422, detail="Could not read DOCX file.")

        elif filename.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(content_bytes))
                slides_out = []
                all_parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    shape_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            shape_texts.append(shape.text.strip())
                    slides_out.append({"number": slide_num, "texts": shape_texts})
                    if shape_texts:
                        all_parts.append(f"[Slide {slide_num}]\n" + "\n".join(shape_texts))
                pptx_text = "\n\n---\n\n".join(all_parts)
                # Return early — skip the global text-length check so image-heavy
                # presentations don't error out; frontend handles empty slides gracefully.
                return {"text": pptx_text[:500_000], "slides": slides_out}
            except HTTPException:
                raise
            except Exception:
                raise HTTPException(status_code=422, detail="Could not read PPTX file.")

        elif filename.endswith((".txt", ".md", ".csv")):
            try:
                text = content_bytes.decode("utf-8", errors="ignore")
            except Exception:
                raise HTTPException(status_code=422, detail="Could not read text file.")

        else:
            raise HTTPException(status_code=400, detail="Unsupported file type. Use PDF, DOCX, PPTX, or TXT.")

        if not text or len(text.strip()) < 10:
            raise HTTPException(status_code=422, detail="No readable text found in this file.")

        return {"text": text[:500_000]}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error in /extract-file-text: {e}")
        raise HTTPException(status_code=500, detail="Failed to extract file text")


@app.get("/")
def health_check():
    """Health check endpoint."""
    return {"status": "ok", "version": "1.0.0"}


@app.get("/domains")
def list_domains():
    """Return available academic domains and their exam modes."""
    from domains import DOMAIN_LIST
    return {"domains": DOMAIN_LIST}


@app.post("/ingest", response_model=IngestResponse)
@limiter.limit("30/minute")
async def ingest(body: IngestRequest, request: Request, authorization: str = Header(default="")):
    """
    Ingest page content for processing.
    Requires authentication. Rate limited.
    """
    try:
        # Require auth to prevent OpenAI credit abuse
        user_id = get_user_id(authorization)
        logger.info(f"Ingesting content for user={user_id[:8]}...")

        # Sanitize inputs
        content = _sanitize_text(body.content, MAX_CONTENT_LENGTH)
        page_url = _validate_url(body.page_url)

        if not content or len(content) < 10:
            raise HTTPException(status_code=400, detail="Content too short")

        logger.info(f"Content length: {len(content)} chars")

        content_id = str(uuid4())

        # Detect if content is from a slideshow
        is_slideshow, slideshow_type = is_slideshow_content(content)

        # Convert image models to dicts for storage
        images_data = []
        if body.images:
            for img in body.images[:10]:  # Max 10 images
                images_data.append({
                    "data": img.data,
                    "slide_index": img.slide_index,
                    "context": img.context,
                    "alt": img.alt,
                })

        # Store content with metadata (keyed by user to prevent cross-user access)
        storage.save_content(
            content_id,
            content,
            page_url,
            metadata={
                "content_type": body.content_type,
                "is_slideshow": is_slideshow,
                "slideshow_type": slideshow_type,
                "user_id": user_id,
                "domain": body.domain,
            },
            images=images_data
        )

        logger.info(f"Ingested content_id={content_id}, slideshow={is_slideshow}, images={len(images_data)}")

        return IngestResponse(
            content_id=content_id,
            content_type=body.content_type,
            detected_slideshow=is_slideshow
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error in /ingest: {e}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to ingest content")


@app.post("/generate", response_model=GenerateResponse)
@limiter.limit("15/minute")
async def generate(body: GenerateRequest, request: Request, authorization: str = Header(default="")):
    """
    Generate study materials from ingested content.
    Requires authentication. Rate limited.
    """
    try:
        user_id = get_user_id(authorization)
        logger.info(f"Generating materials for user={user_id[:8]}...")

        # Enforce freemium usage limit before doing any work
        check_and_increment_usage(user_id)

        content_obj = storage.get_content(body.content_id)
        if not content_obj:
            raise HTTPException(status_code=404, detail="Content not found")

        # Verify content belongs to requesting user
        if content_obj.get("metadata", {}).get("user_id") != user_id:
            raise HTTPException(status_code=403, detail="Access denied")

        raw_text = content_obj["content"]
        metadata = content_obj.get("metadata", {})

        # Check if this is slideshow content
        slide_count = 0
        slides = None
        if metadata.get("is_slideshow"):
            slides = extract_slideshow_content(raw_text)
            if slides:
                # Format all slides as structured XML text — no AI compression
                # so every slide's content reaches generate_study_guide intact
                cleaned = format_slideshow_text(slides)
                slide_count = len(slides)
                logger.info(f"Extracted {slide_count} slides")
            else:
                cleaned = clean_text(raw_text)
        else:
            cleaned = clean_text(raw_text)

        # Chunk the content for processing.
        # When slides are available, chunk_text produces one XML chunk per slide
        # instead of re-splitting the formatted text by paragraphs.
        chunks = chunk_text(cleaned, slides=slides)

        # Analyze images via GPT-4o vision if present, then inject descriptions
        stored_images = content_obj.get("images", [])
        has_images = False
        if stored_images:
            logger.info(f"Analyzing {len(stored_images)} images via vision API...")
            # Build slide text context for vision prompts
            slide_text_map = {}
            if slides:
                for i, s in enumerate(slides):
                    slide_text_map[i] = s.get("content", [])
            image_descs = analyze_images_for_slides(stored_images, slide_text_map)
            if image_descs:
                has_images = True
                logger.info(f"Got {len(image_descs)} image descriptions")
                if slides:
                    slides = inject_image_descriptions(slides, image_descs)
                    cleaned = format_slideshow_text(slides)
                    chunks = chunk_text(cleaned, slides=slides)
                else:
                    cleaned = inject_page_image_descriptions(cleaned, image_descs)
                    chunks = chunk_text(cleaned)
            # Free image data from memory after analysis
            content_obj["images"] = []

        if not chunks:
            logger.warning("No content chunks after cleaning")
            return GenerateResponse(
                notes="No meaningful content found to process.",
                study_guide=None,
                flashcards=None
            )

        # Generate requested materials
        notes_str = None
        study_guide = None
        flashcards = None

        if body.notes:
            logger.info("Generating notes...")
            notes = generate_notes_ai('\n\n'.join(chunks))
            notes_str = '\n'.join(f"- {note}" for note in notes) if notes else "No notes generated."

        if body.study_guide:
            logger.info("Generating study guide...")
            # Use domain from request body, or fall back to what was stored at ingest
            domain = body.domain or metadata.get("domain")
            study_guide = generate_study_guide(chunks, has_images=has_images, domain=domain)

        if body.flashcards:
            logger.info("Generating flashcards...")
            # Derive flashcards from study guide Q&A pairs (1:1) instead of separate AI call
            if study_guide:
                flashcards = []
                lines = study_guide.split('\n')
                current_q = None
                for line in lines:
                    q_match = re.match(r'^Q\d+:\s*(.+)', line)
                    a_match = re.match(r'^A\d+:\s*(.+)', line)
                    if q_match:
                        current_q = q_match.group(1).strip()
                    elif a_match and current_q:
                        flashcards.append({'front': current_q, 'back': a_match.group(1).strip()})
                        current_q = None
                logger.info(f"Created {len(flashcards)} flashcards from study guide Q&A pairs")
            else:
                flashcards = generate_flashcards(chunks)

        return GenerateResponse(
            notes=notes_str,
            study_guide=study_guide,
            flashcards=flashcards
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error in /generate: {e}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to generate materials")


@app.post("/flashcards", response_model=FlashcardResponse)
@limiter.limit("15/minute")
async def create_flashcards(body: FlashcardRequest, request: Request, authorization: str = Header(default="")):
    """
    Generate flashcards from ingested content.
    Requires authentication. Rate limited.
    """
    try:
        user_id = get_user_id(authorization)
        logger.info(f"Generating flashcards for user={user_id[:8]}...")

        # Enforce max_cards limit
        max_cards = min(body.max_cards, 30)

        content_obj = storage.get_content(body.content_id)
        if not content_obj:
            raise HTTPException(status_code=404, detail="Content not found")

        # Verify content belongs to requesting user
        if content_obj.get("metadata", {}).get("user_id") != user_id:
            raise HTTPException(status_code=403, detail="Access denied")

        raw_text = content_obj["content"]
        cleaned = clean_text(raw_text)
        chunks = chunk_text(cleaned)

        if not chunks:
            return FlashcardResponse(flashcards=[], count=0)

        flashcards = generate_flashcards(chunks, max_cards=max_cards)

        return FlashcardResponse(
            flashcards=[{"front": fc["front"], "back": fc["back"]} for fc in flashcards],
            count=len(flashcards)
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error in /flashcards: {e}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to generate flashcards")


@app.post("/chat", response_model=ChatResponse)
@limiter.limit("30/minute")
async def chat(body: ChatRequest, request: Request, authorization: str = Header(default="")):
    """
    Answer questions about the content.
    Requires authentication. Rate limited.
    """
    try:
        user_id = get_user_id(authorization)
        logger.info(f"Chat from user={user_id[:8]}...")

        question = _sanitize_text(body.question, MAX_QUESTION_LENGTH)
        content = _sanitize_text(body.content, MAX_CONTENT_LENGTH)

        if not question:
            raise HTTPException(status_code=400, detail="Question is required")

        if not content:
            return ChatResponse(answer="No content provided. Please capture a page first.")

        # Validate mode
        if body.mode not in ("short", "detailed", "example"):
            raise HTTPException(status_code=400, detail="Invalid mode")

        answer = answer_question(
            question=question,
            context=content,
            mode=body.mode
        )

        return ChatResponse(answer=answer)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error in /chat: {e}\n{traceback.format_exc()}")
        return JSONResponse(
            status_code=500,
            content={"answer": "Error processing your question. Please try again."}
        )

